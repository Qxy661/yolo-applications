"""生成风电叶片缺陷检测PPT — 12页 (中文配图版)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(__file__)
IMG = os.path.join(BASE, 'images')

# 配色
BLUE = RGBColor(0x2E, 0x86, 0xAB)
GREEN = RGBColor(0x28, 0xA7, 0x45)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
RED = RGBColor(0xDC, 0x35, 0x45)
PURPLE = RGBColor(0x6F, 0x42, 0xC1)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
DARK_BLUE = RGBColor(0x1A, 0x56, 0x76)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title(slide, text, left=Inches(0.5), top=Inches(0.2), width=Inches(12), size=Pt(32)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.LEFT


def add_text(slide, text, left, top, width, height, size=Pt(16), color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_bullet(slide, items, left, top, width, height, size=Pt(14), color=BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf


def add_img(slide, fname, left, top, width=None, height=None):
    path = os.path.join(IMG, fname)
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)


# ============================================================
# P1 — 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK_BLUE

add_text(slide, '风电叶片缺陷检测', Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         size=Pt(44), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, '基于YOLO的技术调研与方案设计', Inches(1), Inches(2.7), Inches(11), Inches(0.8),
         size=Pt(28), color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
add_text(slide, 'Technical Survey and Scheme Design for Wind Turbine Blade Defect Detection Based on YOLO',
         Inches(1), Inches(3.5), Inches(11), Inches(0.6), size=Pt(16), color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER)
add_text(slide, '姓名: XXX    学号: XXXXXXXX\n导师: XXX 教授\n自动化科学与电气工程学院\n2026年5月',
         Inches(1), Inches(5), Inches(11), Inches(1.5), size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# P2 — 工业背景与技术挑战
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '风电叶片检测：为什么重要？为什么难？')

# 左图：风电场场景
add_img(slide, 'fig02_wind_farm.png', Inches(0.2), Inches(1.0), width=Inches(6.4))
# 右图：五大挑战
add_img(slide, 'fig03_challenges.png', Inches(6.8), Inches(1.0), width=Inches(6.2))

# 底部要点
add_bullet(slide, [
    '全球风电装机 >1000GW，叶片维修成本占 15-20%',
    '传统人工: 单台 2-4 小时，高空风险，漏检率 15-30%',
    'AI 自动检测: 无人机 15-30 分钟，mAP >90%',
], Inches(0.3), Inches(5.8), Inches(6.2), Inches(1.5), size=Pt(12))


# ============================================================
# P3 — GCB-YOLO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '论文一：GCB-YOLO (2025) — 轻量化+CA+BiFPN')

# 架构图（中文版）
add_img(slide, 'fig04_gcb_yolo.png', Inches(0.3), Inches(1.0), width=Inches(12.5))

add_bullet(slide, [
    '① GhostNet 轻量化 Backbone: 一半正常卷积 + 一半廉价线性变换 → 参数-50%',
    '② CA 坐标注意力: X/Y方向1D池化 → 保留行列坐标，对小目标友好',
    '③ BiFPN 加权特征融合: 双向融合 + 可学习权重 → 多尺度信息更充分',
], Inches(0.5), Inches(4.5), Inches(12), Inches(1.5), size=Pt(13))

# 论文原图：CA注意力对比
add_img(slide, 'ca_attention.png', Inches(0.5), Inches(5.8), width=Inches(4))
add_text(slide, 'CA vs SE vs CBAM 结构对比 (Hou et al., ICCV 2021)',
         Inches(0.5), Inches(7.0), Inches(4), Inches(0.3), size=Pt(9), color=GRAY)

# 论文原图：BiFPN
add_img(slide, 'bifpn.png', Inches(5.0), Inches(5.8), width=Inches(4.5))
add_text(slide, 'BiFPN 双向加权融合 (Tan et al., 2020)',
         Inches(5.0), Inches(7.0), Inches(4.5), Inches(0.3), size=Pt(9), color=GRAY)

# 论文原图：GhostNet
add_img(slide, 'ghostnet_arch.png', Inches(9.8), Inches(5.8), width=Inches(3.3))
add_text(slide, 'GhostNetV1 vs V2 (Han et al., CVPR 2020)',
         Inches(9.8), Inches(7.0), Inches(3.3), Inches(0.3), size=Pt(9), color=GRAY)


# ============================================================
# P4 — WTBD-YOLOv8
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '论文二：WTBD-YOLOv8 (2024) — 更少参数+更高精度')

# 四模块图（中文版）
add_img(slide, 'fig05_wtbd_yolov8.png', Inches(0.3), Inches(1.0), width=Inches(12.5))

add_bullet(slide, [
    '① GhostCBS: 替换标准卷积 → 参数-38.2%',
    '② DFSB-C3k2: 密集特征尺度平衡',
    '③ MHSA-C3k2: 多头自注意力，捕获全局上下文',
    '④ Mini-BiFPN: 2层加权融合 → 参数-50%',
], Inches(0.5), Inches(4.8), Inches(5.5), Inches(2), size=Pt(13))

add_text(slide, 'AP: 98.3% (+2.2%) | 小目标AP: 97.9% (+4.8%) | 参数: 1.99M (-38.2%)',
         Inches(0.5), Inches(6.8), Inches(12), Inches(0.5), size=Pt(14), color=GREEN, bold=True)


# ============================================================
# P5 — LE-YOLO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '论文三：LE-YOLO (2024) — 零参数注意力+高效损失')

# SimAM图（中文版）
add_img(slide, 'fig06_le_yolo.png', Inches(0.3), Inches(1.0), width=Inches(12.5))

add_bullet(slide, [
    '① GSConv: Ghost+Shuffle → 参数-44%，计算-42%',
    '② SimAM 零参数注意力: e_t = (2t-μ)²/(2σ²+ε) → 零参数! 零开销! +2.8% mAP',
    '③ EIoU 损失函数: 分解重叠面积+中心距离+宽高差异 → 比CIoU梯度更精确',
], Inches(0.5), Inches(5.5), Inches(12), Inches(1.5), size=Pt(13))

add_text(slide, '结果: 2.1M参数 | 78.7% mAP | 105.1 FPS    启示: SimAM零参数可免费添加，EIoU直接替代CIoU',
         Inches(0.5), Inches(6.8), Inches(12), Inches(0.5), size=Pt(13), color=ORANGE, bold=True)


# ============================================================
# P6 — 三条技术路线
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '三条技术路线 + 跨论文共同配方')

add_img(slide, 'fig07_three_routes.png', Inches(0.3), Inches(1.0), width=Inches(12.5))

add_bullet(slide, [
    '路线一: 轻量化卷积 (GhostNet/GSConv/GhostCBS) → 参数-38~44%',
    '路线二: 注意力机制 (CA/SimAM/MHSA) → +2~3% mAP，避免CBAM!',
    '路线三: 特征融合 (BiFPN) → +4.6% mAP，小目标AP+4.8%',
], Inches(0.5), Inches(5.0), Inches(12), Inches(1.5), size=Pt(14))

add_text(slide, '共同配方: YOLO + 轻量化 + 注意力 + BiFPN + EIoU/DFL = 更小、更快、更准',
         Inches(0.5), Inches(6.5), Inches(12), Inches(0.5), size=Pt(15), color=ORANGE, bold=True)


# ============================================================
# P7 — 跨领域对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '风电 vs 桥梁：跨领域技术对比')

add_img(slide, 'fig07_cross_domain.png', Inches(0.3), Inches(1.0), width=Inches(12.5))

add_bullet(slide, [
    '共性: 裂纹检测是两个领域的核心需求',
    '共性: 轻量化+注意力+BiFPN是共同最佳实践',
    '启示: 桥梁数据丰富(56K)，可先训练再迁移到风电',
], Inches(0.5), Inches(6.2), Inches(12), Inches(1), size=Pt(14))


# ============================================================
# P8 — 数据集
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '数据集现状与增强策略')

# 数据集柱状图（中文版）
add_img(slide, 'fig08_datasets.png', Inches(0.3), Inches(1.0), width=Inches(6.5))

add_bullet(slide, [
    '数据增强策略:',
    '  Mosaic (4张拼接) → +30% mAP ★★★★★',
    '  MixUp (2张混合) → +10% mAP ★★★★★',
    '  CopyPaste (缺陷复制) → +3.2% mAP',
    '  分辨率渐进 640→1280→1920px ★★★★★',
    '',
    '分辨率渐进训练 (memari-majid, 2024):',
    '  低分辨率收敛快 → 高分辨率保留小缺陷',
], Inches(7.0), Inches(1.2), Inches(6), Inches(5), size=Pt(13))


# ============================================================
# P9 — 损失函数
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '损失函数演进')

# 损失函数演进图（中文版）
add_img(slide, 'fig09_loss.png', Inches(0.3), Inches(1.0), width=Inches(12.5))

add_bullet(slide, [
    '每步改进: +重叠度 → +包围框 → +中心距离 → +宽高比 → +分解 → +分布',
    '',
    '风电推荐: EIoU + DFL',
    '  EIoU: 细长裂纹的宽高比梯度更稳定',
    '  DFL: 5px裂纹偏移2px=40%误差，分布建模更鲁棒',
], Inches(0.5), Inches(5.0), Inches(12), Inches(2), size=Pt(13))


# ============================================================
# P10 — 推理链
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '从调研到方案：推理链')

add_img(slide, 'fig10_reasoning.png', Inches(1.0), Inches(1.0), width=Inches(11))

add_bullet(slide, [
    '论文结论 → GhostNet+CA+BiFPN验证有效',
    'YOLOv11选择 → C3k2+C2PSA+解耦头+DFL',
    '方案确定 → YOLOv11+GhostNet+CA+BiFPN+EIoU/DFL+SAHI',
], Inches(0.5), Inches(6.2), Inches(12), Inches(1), size=Pt(14))


# ============================================================
# P11 — 方案架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '拟采用方案总体架构')

add_img(slide, 'fig11_architecture.png', Inches(0.3), Inches(1.0), width=Inches(12.5))

add_bullet(slide, [
    'YOLOv11 + GhostNet + CA + BiFPN + EIoU/DFL',
    '数据增强: Mosaic+MixUp+CopyPaste+分辨率渐进',
    '推理优化: SAHI切片 (640×640, overlap=0.2)',
    '预期: mAP>90% | <10MB | >60FPS',
], Inches(0.5), Inches(6.0), Inches(12), Inches(1.2), size=Pt(14))


# ============================================================
# P12 — 总结+致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = DARK_BLUE

add_text(slide, '调研总结', Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         size=Pt(36), color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_bullet(slide, [
    '1. 精读3篇核心论文: GCB-YOLO / WTBD-YOLOv8 / LE-YOLO',
    '2. 归纳三条路线: 轻量化 → 注意力 → 特征融合',
    '3. 跨领域对比: 风电+桥梁双领域分析',
    '4. 关键发现: GhostNet参数-38%, CA优于CBAM, BiFPN小目标+4.8%',
    '5. 方案: YOLOv11 + GhostNet + CA + BiFPN + EIoU/DFL + SAHI',
    '6. 预期: mAP>90%, <10MB, >60FPS',
], Inches(1), Inches(1.5), Inches(11), Inches(3.5), size=Pt(16), color=WHITE)

add_text(slide, '参考文献', Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         size=Pt(18), color=RGBColor(0xBB, 0xDE, 0xFB), bold=True)

refs = [
    '[1] Zhang et al., GCB-YOLO, Wind Energy 2025',
    '[2] Tong et al., WTBD-YOLOv8, Sustainability 2024',
    '[3] Fu et al., LE-YOLO, IEEE Access 2024',
    '[4] Hou et al., Coordinate Attention, ICCV 2021',
    '[5] Han et al., GhostNet, CVPR 2020',
    '[6] Ultralytics, YOLOv11, GitHub 2024',
    '[7] Memari et al., YOLO for Wind Turbine: A Review, RSER 2024',
]
add_bullet(slide, refs, Inches(1), Inches(5.0), Inches(11), Inches(2), size=Pt(11), color=RGBColor(0x90, 0xCA, 0xF9))

add_text(slide, '感谢聆听！恳请各位老师批评指正', Inches(1), Inches(6.8), Inches(11), Inches(0.5),
         size=Pt(20), color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# 保存
# ============================================================
out = os.path.join(BASE, 'wind_turbine_yolo_ppt.pptx')
prs.save(out)
print(f'PPT saved: {out}')
print(f'Total slides: {len(prs.slides)}')
